#!/usr/bin/env python3
"""Unified file generator for chat tool calls.

Reads a JSON spec from stdin and writes a binary file to argv[1].
Supported kinds: docx, xlsx, pdf, pptx.

Spec format:
  {
    "kind": "docx",
    "filename": "report.docx",
    "spec": { ...kind-specific schema... }
  }

Exit codes: 0 = success, 1 = bad input, 2 = library missing, 3 = generation error.
Outputs JSON to stdout: { ok: bool, kind, filename, sizeBytes, error? }
"""
from __future__ import annotations
import json
import os
import sys
from pathlib import Path


# ─────────────────────────────────────────────────────────────────────────────
# DOCX — uses python-docx
# spec: { title?, sections: [{heading?, level?, paragraphs?:[str], lists?:[{type:'bullet'|'number', items:[str]}], table?:{headers:[], rows:[[]]} }] }
# ─────────────────────────────────────────────────────────────────────────────
def gen_docx(out_path: str, spec: dict) -> None:
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    if spec.get("title"):
        doc.add_heading(spec["title"], level=0)

    for sec in spec.get("sections", []):
        if sec.get("heading"):
            doc.add_heading(sec["heading"], level=int(sec.get("level", 1)))
        for para in sec.get("paragraphs", []) or []:
            p = doc.add_paragraph(para)
        for lst in sec.get("lists", []) or []:
            style = "List Bullet" if lst.get("type", "bullet") == "bullet" else "List Number"
            for item in lst.get("items", []):
                doc.add_paragraph(item, style=style)
        tbl = sec.get("table")
        if tbl and tbl.get("rows"):
            headers = tbl.get("headers", [])
            rows = tbl["rows"]
            ncols = max(len(headers), max((len(r) for r in rows), default=0))
            t = doc.add_table(rows=1 + len(rows), cols=ncols)
            t.style = "Light Grid Accent 1"
            for i, h in enumerate(headers):
                t.rows[0].cells[i].text = str(h)
            for ri, row in enumerate(rows):
                for ci, val in enumerate(row):
                    t.rows[1 + ri].cells[ci].text = str(val)

    doc.save(out_path)


# ─────────────────────────────────────────────────────────────────────────────
# XLSX — uses openpyxl
# spec: { sheets: [{name?, headers?:[str], rows:[[any]] }] }
# ─────────────────────────────────────────────────────────────────────────────
def gen_xlsx(out_path: str, spec: dict) -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment

    wb = Workbook()
    wb.remove(wb.active)

    sheets = spec.get("sheets") or [{"name": "Sheet1", "rows": []}]
    for i, sh in enumerate(sheets):
        ws = wb.create_sheet(title=(sh.get("name") or f"Sheet{i+1}")[:31])
        headers = sh.get("headers") or []
        rows    = sh.get("rows")    or []
        if headers:
            ws.append(headers)
            bold = Font(bold=True, color="FFFFFFFF")
            fill = PatternFill(start_color="FF305496", end_color="FF305496", fill_type="solid")
            for cell in ws[1]:
                cell.font = bold
                cell.fill = fill
                cell.alignment = Alignment(horizontal="center")
            ws.freeze_panes = "A2"
        for row in rows:
            ws.append(row)
        # Autosize-ish: cap to 50
        for col_cells in ws.columns:
            max_len = max((len(str(c.value)) for c in col_cells if c.value is not None), default=0)
            ws.column_dimensions[col_cells[0].column_letter].width = min(max_len + 2, 50)

    wb.save(out_path)


# ─────────────────────────────────────────────────────────────────────────────
# PDF — uses reportlab
# spec: { title?, content_markdown?, sections?:[{heading?, paragraphs?:[str], lists?:[{items:[str]}], table?:{headers, rows}}] }
# ─────────────────────────────────────────────────────────────────────────────
def gen_pdf(out_path: str, spec: dict) -> None:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib import colors
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                    Table, TableStyle, ListFlowable, ListItem)
    from reportlab.lib.units import cm

    doc = SimpleDocTemplate(out_path, pagesize=A4,
                            leftMargin=2*cm, rightMargin=2*cm,
                            topMargin=2*cm, bottomMargin=2*cm)
    styles = getSampleStyleSheet()
    story = []
    if spec.get("title"):
        story.append(Paragraph(spec["title"], styles["Title"]))
        story.append(Spacer(1, 0.5*cm))

    # Simple markdown body if provided
    md = spec.get("content_markdown")
    if md:
        for line in md.split("\n"):
            line = line.rstrip()
            if not line:
                story.append(Spacer(1, 0.3*cm))
                continue
            if line.startswith("# "):
                story.append(Paragraph(line[2:], styles["Heading1"]))
            elif line.startswith("## "):
                story.append(Paragraph(line[3:], styles["Heading2"]))
            elif line.startswith("### "):
                story.append(Paragraph(line[4:], styles["Heading3"]))
            elif line.startswith("- ") or line.startswith("* "):
                story.append(Paragraph("• " + line[2:], styles["BodyText"]))
            else:
                story.append(Paragraph(line, styles["BodyText"]))

    for sec in spec.get("sections", []) or []:
        if sec.get("heading"):
            story.append(Paragraph(sec["heading"], styles["Heading1"]))
        for p in sec.get("paragraphs", []) or []:
            story.append(Paragraph(p, styles["BodyText"]))
        for lst in sec.get("lists", []) or []:
            story.append(ListFlowable(
                [ListItem(Paragraph(it, styles["BodyText"])) for it in lst.get("items", [])],
                bulletType="bullet"))
        tbl = sec.get("table")
        if tbl and tbl.get("rows"):
            data = [tbl.get("headers", [])] + tbl["rows"]
            t = Table(data, repeatRows=1)
            t.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#305496")),
                ("TEXTCOLOR",  (0, 0), (-1, 0), colors.white),
                ("FONTNAME",   (0, 0), (-1, 0), "Helvetica-Bold"),
                ("GRID",       (0, 0), (-1, -1), 0.25, colors.grey),
                ("VALIGN",     (0, 0), (-1, -1), "TOP"),
            ]))
            story.append(t)
        story.append(Spacer(1, 0.4*cm))

    doc.build(story)


# ─────────────────────────────────────────────────────────────────────────────
# PPTX — uses python-pptx
# spec: { slides: [{title?, bullets?:[str], body?, notes?}] }
# ─────────────────────────────────────────────────────────────────────────────
def gen_pptx(out_path: str, spec: dict) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    title_only = prs.slide_layouts[5]   # Title Only
    content    = prs.slide_layouts[1]   # Title + Content

    for sl in spec.get("slides", []) or []:
        bullets = sl.get("bullets") or []
        body    = sl.get("body")
        if bullets or body:
            slide = prs.slides.add_slide(content)
            if slide.shapes.title:
                slide.shapes.title.text = sl.get("title", "")
            tf = slide.placeholders[1].text_frame
            if bullets:
                first = True
                for b in bullets:
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    p.text = str(b)
                    p.level = 0
                    first = False
            elif body:
                tf.text = str(body)
        else:
            slide = prs.slides.add_slide(title_only)
            if slide.shapes.title:
                slide.shapes.title.text = sl.get("title", "")

        if sl.get("notes"):
            slide.notes_slide.notes_text_frame.text = str(sl["notes"])

    prs.save(out_path)


# ─────────────────────────────────────────────────────────────────────────────
# Entry point
# ─────────────────────────────────────────────────────────────────────────────
GENERATORS = {"docx": gen_docx, "xlsx": gen_xlsx, "pdf": gen_pdf, "pptx": gen_pptx}


def main():
    if len(sys.argv) < 2:
        print(json.dumps({"ok": False, "error": "usage: file-gen.py <output-path>"}))
        sys.exit(1)

    out_path = sys.argv[1]
    try:
        raw = sys.stdin.read()
        req = json.loads(raw)
    except Exception as e:
        print(json.dumps({"ok": False, "error": f"invalid JSON input: {e}"}))
        sys.exit(1)

    kind = (req.get("kind") or "").lower()
    if kind not in GENERATORS:
        print(json.dumps({"ok": False, "error": f"unsupported kind: {kind}"}))
        sys.exit(1)

    spec = req.get("spec") or {}
    try:
        os.makedirs(os.path.dirname(out_path), exist_ok=True)
        GENERATORS[kind](out_path, spec)
    except ImportError as e:
        print(json.dumps({"ok": False, "error": f"missing python library: {e}"}))
        sys.exit(2)
    except Exception as e:
        print(json.dumps({"ok": False, "error": f"generation failed: {type(e).__name__}: {e}"}))
        sys.exit(3)

    size = os.path.getsize(out_path)
    print(json.dumps({"ok": True, "kind": kind, "path": out_path, "sizeBytes": size}))
    sys.exit(0)


if __name__ == "__main__":
    main()
